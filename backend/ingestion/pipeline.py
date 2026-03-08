"""
Document ingestion pipeline — handles text extraction, chunking, and embedding storage.
Supports: PDF, DOCX, TXT, CSV, PPTX
"""
import os
from typing import List, Dict, Any

# Text Extractors
def extract_pdf(file_path: str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)

def extract_docx(file_path: str) -> str:
    import docx2txt
    return docx2txt.process(file_path)

def extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_csv(file_path: str) -> str:
    import csv
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)

def extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n\n".join(texts)

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".doc": extract_docx,
    ".txt": extract_txt,
    ".csv": extract_csv,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
}

def extract_text(file_path: str) -> str:
    """Extract text from a document based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractor(file_path)

def chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks for embedding."""
    if not text:
        return []
    
    chunks = []
    # Split by double newline first (paragraphs)
    paragraphs = text.split('\n\n')
    
    current_chunk = ""
    
    for para in paragraphs:
        # If the paragraph itself is larger than chunk_size, we should split it further
        if len(para) > chunk_size:
            # Add existing chunk if have one
            if current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = ""
            
            # Simple naive chop for massive paragraphs
            for i in range(0, len(para), chunk_size - chunk_overlap):
                chunks.append(para[i:i + chunk_size])
            continue
            
        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk += ("\n\n" if current_chunk else "") + para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            
            # Start new chunk with some overlap from previous if possible
            if current_chunk and chunk_overlap > 0:
                overlap_text = current_chunk[-chunk_overlap:]
                # Try to cut at a clean boundary
                idx = overlap_text.find('. ')
                if idx != -1:
                    overlap_text = overlap_text[idx+2:]
                current_chunk = overlap_text + "\n\n" + para
            else:
                current_chunk = para
                
    if current_chunk:
        chunks.append(current_chunk.strip())
        
    return chunks

def process_document(file_path: str, file_id: str) -> Dict[str, Any]:
    """Full ingestion: extract text → chunk → return metadata."""
    text = extract_text(file_path)
    chunks = chunk_text(text)
    filename = os.path.basename(file_path)
    return {
        "file_id": file_id,
        "filename": filename,
        "text": text,
        "chunks": chunks,
        "num_chunks": len(chunks),
        "file_size": os.path.getsize(file_path),
    }
